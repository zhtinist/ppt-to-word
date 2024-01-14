from pptx import Presentation
from docx import Document
from docx.shared import Pt
import os
from os.path import join, isfile, isdir
from urllib.parse import unquote
from pathlib import Path
import zipfile
import tempfile
import shutil

def is_valid_pptx_file(file_path):
    temp_dir = tempfile.mkdtemp()  # 创建临时目录
    try:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        return True
    except zipfile.BadZipFile:
        return False
    finally:
        shutil.rmtree(temp_dir)  # 自动清理临时目录

def extract_text_from_ppt(ppt_file):
    temp_dir = tempfile.mkdtemp()  # 创建临时目录
    try:
        if not is_valid_pptx_file(ppt_file):
            return ""

        presentation = Presentation(ppt_file)
        text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text
    finally:
        shutil.rmtree(temp_dir)  # 自动清理临时目录

def create_word_document(ppt_path, output_folder):
    ppt_filename = os.path.basename(ppt_path)
    ppt_name, _ = os.path.splitext(ppt_filename)

    word_filename = f"{ppt_name}.docx"
    word_path = join(output_folder, word_filename)

    ppt_text = extract_text_from_ppt(ppt_path)

    document = Document()
    document.styles['Normal'].font.name = '微软雅黑'
    
    ppt_text = ppt_text.encode('utf-8').decode('utf-8')
    
    paragraph = document.add_paragraph(ppt_text)

    for run in paragraph.runs:
        run.font.size = Pt(11)

    document.save(word_path)
    print(f"已处理：{ppt_filename} -> {word_filename}")

def process_ppts(path):
    path = unquote(path)  # 解码中文路径
    path = Path(path)  # 使用pathlib规范化路径
    if path.is_dir():
        for file_path in path.glob("*.ppt*"):
            if isfile(file_path) and is_valid_pptx_file(file_path):
                create_word_document(file_path, str(path))
    elif path.is_file() and path.suffix.lower() in (".ppt", ".pptx") and is_valid_pptx_file(str(path)):
        create_word_document(str(path), str(path.parent))
    else:
        print(f"错误：输入的路径既不是有效的文件夹，也不是有效的 PowerPoint 文件。")

if __name__ == "__main__":
    while True:
        input_path = input("请输入待处理文件夹或 PowerPoint 文件的路径：")
        process_ppts(input_path)

        user_input = input("按任意键继续，或输入 'q' 退出程序: ")
        if user_input.lower() == 'q':
            break
